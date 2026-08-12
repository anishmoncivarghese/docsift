"""Regenerate `deck.pptx`, the fixture that proves slide attribution.

Run with `uv run python tests/fixtures/make_deck.py`. The output is committed;
this exists so the fixture can be explained and rebuilt.

60 slides because the point of retrieval is a deck too long to read whole, and
because a one-slide deck would pass an attribution test by accident.
"""

from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).parent
SLIDE_COUNT = 60
# Only slide 14 mentions this. A search for it must come back citing 14.
NEEDLE = "Vendor dependency remains unresolved in the APAC corridor"


def main() -> None:
    prs = Presentation()
    for number in range(1, SLIDE_COUNT + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Section {number}"
        body = slide.placeholders[1].text_frame
        body.text = NEEDLE if number == 14 else f"Routine content for slide {number}."
        body.add_paragraph().text = f"Owner: team {number % 7}"
    prs.save(HERE / "deck.pptx")
    print("wrote", HERE / "deck.pptx")


if __name__ == "__main__":
    main()
